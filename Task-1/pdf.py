import os
import io
import numpy as np
import pandas as pd
import fitz  
import pytesseract
import cv2
import matplotlib.pyplot as plt
from pptx import Presentation
from PIL import Image
from sentence_transformers import SentenceTransformer
import faiss
from transformers import pipeline

class DocumentRAGPipeline:
    def __init__(self, file_path):
        self.file_path = file_path
        self.file_extension = os.path.splitext(file_path)[1].lower()
        
        # Embedding model
        self.embedding_model = SentenceTransformer('all-MiniLM-L6-v2')
        
        # Vector storage
        self.vector_dimension = self.embedding_model.get_sentence_embedding_dimension()
        self.index = faiss.IndexFlatL2(self.vector_dimension)
        
        # Document storage
        self.text_chunks = []
        self.metadata = []
        
        # Local language model for response generation
        self.llm = pipeline("text-generation", model="gpt2")

    def extract_document_content(self):
        if self.file_extension == '.pptx':
            return self._extract_ppt_content()
        elif self.file_extension == '.pdf':
            return self._extract_pdf_content()
        else:
            raise ValueError(f"Unsupported file type: {self.file_extension}")

    def _extract_ppt_content(self):
        slide_pages = []
        prs = Presentation(self.file_path)
        
        for slide_num, slide in enumerate(prs.slides):
            slide_page = {
                "slide_num": slide_num,
                "text": [],
                "images": [],
                "tables": []
            }

            # Extract text
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        slide_page["text"].append(paragraph.text)

            # Extract tables
            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                    rows = [[cell.text for cell in row.cells] for row in table.rows]
                    df = pd.DataFrame(rows[1:], columns=rows[0])
                    slide_page["tables"].append(df)

            # Extract images with OCR
            for shape in slide.shapes:
                if "Picture" in shape.name:
                    try:
                        image_stream = shape.image.blob
                        image = Image.open(io.BytesIO(image_stream))
                        ocr_text = pytesseract.image_to_string(image)
                        slide_page["images"].append(ocr_text)
                    except Exception as e:
                        print(f"Error extracting image: {e}")

            slide_pages.append(slide_page)
        
        return slide_pages

    def _extract_pdf_content(self):
        pdf_pages = []
        doc = fitz.open(self.file_path)
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            page_content = {
                "page_num": page_num,
                "text": page.get_text(),
                "images": []
            }
            
            # Extract images with OCR
            for img_index, img in enumerate(page.get_images()):
                try:
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    
                    # Convert bytes to PIL Image
                    image = Image.open(io.BytesIO(image_bytes))
                    ocr_text = pytesseract.image_to_string(image)
                    page_content["images"].append(ocr_text)
                except Exception as e:
                    print(f"Error extracting PDF image: {e}")
            
            pdf_pages.append(page_content)
        
        return pdf_pages

    def preprocess_document(self, content):
        self.text_chunks.clear()
        self.metadata.clear()

        for item in content:
            # Extract text from different sources
            text_sources = []
            if isinstance(item, dict):
                text_sources.extend(item.get('text', []))
                text_sources.extend(item.get('images', []))
                
                # Add table contents as text
                for table in item.get('tables', []):
                    text_sources.append(table.to_string(index=False))
            
            # Split and add text chunks
            for source_text in text_sources:
                chunks = self._split_text(source_text)
                for chunk in chunks:
                    self.text_chunks.append(chunk)
                    self.metadata.append({
                        'source': self.file_path,
                        'page_number': item.get('slide_num', item.get('page_num', 0))
                    })

    def _split_text(self, text, chunk_size=200, overlap=50):
        words = text.split()
        chunks = []
        
        for i in range(0, len(words), chunk_size - overlap):
            chunk = ' '.join(words[i:i+chunk_size])
            chunks.append(chunk)
        
        return chunks

    def create_vector_index(self):
        if not self.text_chunks:
            raise ValueError("No text chunks available. Call preprocess_document first.")
        
        embeddings = self.embedding_model.encode(self.text_chunks)
        self.index.add(embeddings)

    def semantic_search(self, query, top_k=3):
        query_embedding = self.embedding_model.encode([query])
        distances, indices = self.index.search(query_embedding, top_k)
        
        results = []
        for idx in indices[0]:
            results.append({
                'text': self.text_chunks[idx],
                'metadata': self.metadata[idx]
            })
        
        return results

    def generate_response(self, query, context):
        context_text = " ".join([item['text'] for item in context])
        prompt = f"Answer the query: {query}\n\nContext: {context_text}"
        
        response = self.llm(prompt, max_new_tokens=150, num_return_sequences=1)[0]['generated_text']
        return response

    def comparison_query(self, comparison_terms, top_k=3):
        comparison_results = []
        
        for term in comparison_terms:
            search_results = self.semantic_search(term, top_k)
            
            for result in search_results:
                comparison_results.append({
                    'term': term,
                    'text': result['text'],
                    'page': result['metadata']['page_number']
                })
        
        return pd.DataFrame(comparison_results)

def main():
    file_path = 'test.pptx'  # or 'test.pdf'
    pipeline = DocumentRAGPipeline(file_path)
    document_content = pipeline.extract_document_content()
    pipeline.preprocess_document(document_content)
    pipeline.create_vector_index()
    query = "unemployment rates by education level"
    semantic_results = pipeline.semantic_search(query)
    response = pipeline.generate_response(query, semantic_results)
    print("Query Response:", response)
    
    # Comparison query
    comparison_terms = ['education', 'unemployment']
    comparison_results = pipeline.comparison_query(comparison_terms)
    print("\nComparison Results:")
    print(comparison_results)
    page_6_table = document_content[5]["tables"]  # Assuming the table is on page 6 (index 5)
    if page_6_table:
        print("\nPage 6 Table:")
        print(page_6_table[0].to_string(index=False))

if __name__ == "__main__":
    main()